#!/usr/bin/env python3
import time
import base64
import logging
import os
from urllib.parse import urlparse, parse_qs
import argparse

import undetected_chromedriver as uc
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

# Import library anti-captcha
from python_anticaptcha import AnticaptchaClient, NoCaptchaTaskProxylessTask

# Konfigurasi logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S"
)
logger = logging.getLogger(__name__)

# Konfigurasi default
CAPTCHA_SERVICE_API_KEY = os.getenv("CAPTCHA_SERVICE_API_KEY", "your_captcha_service_api_key")
DEFAULT_WEBSITE_URL = "https://target-website.com"


def extract_sitekey(iframe_element):
    """
    Ekstrak sitekey dari elemen iframe CAPTCHA.
    """
    sitekey = iframe_element.get_attribute("data-sitekey")
    if sitekey:
        return sitekey
    src = iframe_element.get_attribute("src")
    if src:
        parsed = urlparse(src)
        params = parse_qs(parsed.query)
        return params.get("k", [None])[0]
    return None


def solve_captcha(driver, website_url, timeout=30):
    """
    Cari dan selesaikan CAPTCHA menggunakan layanan AntiCaptcha.
    """
    logger.info("Mencari elemen iframe CAPTCHA...")
    try:
        iframe = WebDriverWait(driver, timeout).until(
            EC.presence_of_element_located((By.XPATH, "//iframe[contains(@src, 'recaptcha')]"))
        )
    except Exception as e:
        logger.error("Elemen CAPTCHA tidak ditemukan: %s", e)
        return None

    sitekey = extract_sitekey(iframe)
    if not sitekey:
        logger.error("Sitekey tidak dapat diekstrak dari iframe CAPTCHA.")
        return None

    logger.info("Sitekey CAPTCHA ditemukan: %s", sitekey)
    try:
        client = AnticaptchaClient(CAPTCHA_SERVICE_API_KEY)
        task = NoCaptchaTaskProxylessTask(
            website_url=website_url,
            website_key=sitekey
        )
        job = client.createTask(task)
        logger.info("Tugas CAPTCHA dibuat. Menunggu solusi...")
        job.join()
        solution = job.get_solution_response()
        logger.info("CAPTCHA berhasil diselesaikan.")
        return solution
    except Exception as e:
        logger.error("Gagal menyelesaikan CAPTCHA: %s", e)
        return None


def inject_captcha_solution(driver, solution):
    """
    Injeksi solusi CAPTCHA ke dalam elemen yang sesuai.
    """
    try:
        driver.execute_script(
            "var elem = document.getElementById('g-recaptcha-response');"
            "if(elem){ elem.style.display = 'block'; elem.value = arguments[0]; }", solution
        )
        logger.info("Solusi CAPTCHA telah diinjeksikan.")
    except Exception as e:
        logger.error("Gagal menginjeksi solusi CAPTCHA: %s", e)


def submit_captcha_form(driver):
    """
    Submit form CAPTCHA. Jika form dengan id khusus tidak ditemukan, submit form pertama.
    """
    try:
        form = driver.find_element(By.XPATH, "//form[contains(@id, 'recaptcha')]")
        form.submit()
        logger.info("Form CAPTCHA disubmit.")
    except Exception as e:
        logger.warning("Form dengan id 'recaptcha' tidak ditemukan, mencoba submit form pertama.")
        try:
            driver.execute_script("document.forms[0].submit();")
        except Exception as ex:
            logger.error("Gagal submit form: %s", ex)


def save_pdf_from_chrome(driver, output_path):
    """
    Simpan halaman saat ini ke PDF menggunakan perintah Chrome DevTools Protocol.
    """
    logger.info("Menyimpan halaman sebagai PDF...")
    try:
        pdf = driver.execute_cdp_cmd("Page.printToPDF", {"printBackground": True})
        with open(output_path, "wb") as f:
            f.write(base64.b64decode(pdf["data"]))
        logger.info("PDF berhasil disimpan ke %s", output_path)
    except Exception as e:
        logger.error("Gagal menyimpan PDF: %s", e)


def convert_to_docx(html_content, output_path):
    """
    Konversi konten HTML ke DOCX dengan mengekstrak teks menggunakan BeautifulSoup.
    """
    logger.info("Mengonversi halaman ke DOCX...")
    try:
        from bs4 import BeautifulSoup
        from docx import Document

        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text(separator="\n")
        doc = Document()
        doc.add_paragraph(text)
        doc.save(output_path)
        logger.info("File DOCX berhasil disimpan ke %s", output_path)
    except Exception as e:
        logger.error("Gagal mengonversi ke DOCX: %s", e)


def convert_to_pptx(html_content, output_path):
    """
    Konversi konten HTML ke PPTX dengan mengekstrak teks dan memasukkannya ke slide.
    """
    logger.info("Mengonversi halaman ke PPTX...")
    try:
        from bs4 import BeautifulSoup
        from pptx import Presentation
        from pptx.util import Inches

        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text(separator="\n")
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Menggunakan layout kosong
        slide = prs.slides.add_slide(slide_layout)
        left = top = Inches(1)
        width = Inches(8)
        height = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = text
        prs.save(output_path)
        logger.info("File PPTX berhasil disimpan ke %s", output_path)
    except Exception as e:
        logger.error("Gagal mengonversi ke PPTX: %s", e)


def main():
    parser = argparse.ArgumentParser(
        description="Mengonversi website ke format PDF, DOCX, atau PPTX."
    )
    parser.add_argument(
        "--format",
        type=str,
        choices=["pdf", "docx", "pptx"],
        default="pdf",
        help="Format output (pdf, docx, atau pptx)"
    )
    parser.add_argument(
        "--website",
        type=str,
        default=DEFAULT_WEBSITE_URL,
        help="URL website yang akan dikonversi"
    )
    parser.add_argument(
        "--output",
        type=str,
        help="Nama file output (opsional)"
    )
    args = parser.parse_args()

    output_format = args.format.lower()
    output_file = args.output if args.output else f"output.{output_format}"
    website_url = args.website

    logger.info("Memulai konversi untuk website: %s", website_url)
    logger.info("Format output yang dipilih: %s", output_format)

    options = uc.ChromeOptions()
    options.add_argument("--headless=new")
    options.add_argument("--disable-blink-features=AutomationControlled")

    driver = uc.Chrome(options=options)
    driver.set_page_load_timeout(60)

    try:
        logger.info("Mengakses %s", website_url)
        driver.get(website_url)

        # Tangani CAPTCHA jika ada
        try:
            captcha_iframe = WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.XPATH, "//iframe[contains(@src, 'recaptcha')]"))
            )
            logger.info("CAPTCHA terdeteksi di halaman.")
            captcha_solution = solve_captcha(driver, website_url)
            if captcha_solution:
                inject_captcha_solution(driver, captcha_solution)
                submit_captcha_form(driver)
                time.sleep(5)
            else:
                logger.error("Tidak mendapatkan solusi CAPTCHA. Proses dihentikan.")
                return
        except Exception as e:
            logger.info("Tidak ada CAPTCHA yang terdeteksi atau terjadi error: %s", e)

        # Tunggu hingga halaman termuat dengan sempurna
        WebDriverWait(driver, 30).until(
            EC.presence_of_element_located((By.TAG_NAME, "body"))
        )
        logger.info("Halaman termuat dengan sempurna.")

        # Ambil konten HTML halaman
        html_content = driver.page_source

        # Lakukan konversi sesuai format yang dipilih
        if output_format == "pdf":
            save_pdf_from_chrome(driver, output_file)
        elif output_format == "docx":
            convert_to_docx(html_content, output_file)
        elif output_format == "pptx":
            convert_to_pptx(html_content, output_file)
        else:
            logger.error("Format output tidak dikenal: %s", output_format)

    except Exception as ex:
        logger.error("Terjadi kesalahan selama proses konversi: %s", ex)
    finally:
        driver.quit()
        logger.info("Driver telah ditutup.")


if __name__ == "__main__":
    main()
